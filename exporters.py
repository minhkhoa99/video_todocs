"""
Multi-format exporters: docx, pdf, xlsx, pptx, md, json
"""
import json
import platform
from pathlib import Path


def export_all(data: dict, output_dir: Path, base_name: str, formats: list) -> list:
    output_files = []
    for fmt in formats:
        try:
            if fmt == "json":
                f = export_json(data, output_dir, base_name)
            elif fmt == "md":
                f = export_markdown(data, output_dir, base_name)
            elif fmt == "docx":
                f = export_docx(data, output_dir, base_name)
            elif fmt == "pdf":
                f = export_pdf(data, output_dir, base_name)
            elif fmt == "xlsx":
                f = export_xlsx(data, output_dir, base_name)
            elif fmt == "pptx":
                f = export_pptx(data, output_dir, base_name)
            else:
                print(f"   ⚠️  Bỏ qua format không hỗ trợ: {fmt}")
                continue
            output_files.append(f)
            print(f"   ✓ {fmt.upper():5} → {f}")
        except Exception as e:
            print(f"   ✗ {fmt.upper():5} thất bại: {e}")
            import traceback
            traceback.print_exc()
    return output_files


# ----------------------- JSON -----------------------
def export_json(data, output_dir, base_name):
    path = output_dir / f"{base_name}.json"
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return path


# --------------------- Markdown ---------------------
def export_markdown(data, output_dir, base_name):
    path = output_dir / f"{base_name}.md"
    md = []
    md.append(f"# {data.get('title', base_name)}\n")
    md.append(f"- **Loại nội dung:** {data.get('content_type', 'N/A')}")
    md.append(f"- **Ngôn ngữ:** {data.get('language', 'N/A')}")
    md.append(f"- **Thời lượng ước tính:** {data.get('duration_estimate', 'N/A')}\n")

    if data.get("summary"):
        md.append("## Tóm tắt\n")
        md.append(data["summary"] + "\n")

    sections = data.get("sections", [])
    if sections:
        md.append("## Nội dung chi tiết\n")
        for i, s in enumerate(sections, 1):
            md.append(f"### {i}. [{s.get('timestamp_start', '')} - {s.get('timestamp_end', '')}] {s.get('title', '')}\n")
            if s.get("visual_content"):
                md.append(f"**📺 Hình ảnh/Slide:**\n\n{s['visual_content']}\n")
            if s.get("spoken_content"):
                md.append(f"**🎤 Lời nói:**\n\n{s['spoken_content']}\n")
            if s.get("key_points"):
                md.append("**🔑 Ý chính:**\n")
                for kp in s["key_points"]:
                    md.append(f"- {kp}")
                md.append("")
            if s.get("notes"):
                md.append(f"**📝 Ghi chú:** {s['notes']}\n")

    exercises = data.get("exercises", [])
    if exercises:
        md.append("## Bài tập\n")
        for ex in exercises:
            md.append(f"### {ex.get('number', '')} {('— ' + ex['topic']) if ex.get('topic') else ''}\n")
            if ex.get("timestamp"):
                md.append(f"*Thời điểm: {ex['timestamp']}*\n")
            md.append(f"**Đề bài:**\n\n{ex.get('problem_statement', '')}\n")
            if ex.get("given_info"):
                md.append(f"**Dữ kiện:**\n\n{ex['given_info']}\n")
            if ex.get("diagram_description"):
                md.append(f"**Hình vẽ:**\n\n{ex['diagram_description']}\n")
            if ex.get("solution"):
                md.append(f"**Lời giải:**\n\n{ex['solution']}\n")
            if ex.get("answer"):
                md.append(f"**✓ Đáp án:** {ex['answer']}\n")

    key_terms = data.get("key_terms", [])
    if key_terms:
        md.append("## Thuật ngữ & Định nghĩa\n")
        for k in key_terms:
            md.append(f"- **{k.get('term', '')}**: {k.get('definition', '')}")
        md.append("")

    if data.get("full_transcript"):
        md.append("## Toàn bộ lời nói (transcript)\n")
        md.append(data["full_transcript"])

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(md))
    return path


# ----------------------- DOCX -----------------------
def export_docx(data, output_dir, base_name):
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    path = output_dir / f"{base_name}.docx"
    doc = Document()

    # Default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Title
    title = doc.add_heading(data.get("title", base_name), level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(f"Loại: {data.get('content_type', 'N/A')}  |  ").italic = True
    p.add_run(f"Ngôn ngữ: {data.get('language', 'N/A')}  |  ").italic = True
    p.add_run(f"Thời lượng: {data.get('duration_estimate', 'N/A')}").italic = True

    if data.get("summary"):
        doc.add_heading("Tóm tắt", level=1)
        doc.add_paragraph(data["summary"])

    sections = data.get("sections", [])
    if sections:
        doc.add_heading("Nội dung chi tiết", level=1)
        for i, s in enumerate(sections, 1):
            doc.add_heading(
                f"{i}. [{s.get('timestamp_start', '')} - {s.get('timestamp_end', '')}] {s.get('title', '')}",
                level=2,
            )
            if s.get("visual_content"):
                p = doc.add_paragraph()
                p.add_run("📺 Hình ảnh/Slide: ").bold = True
                p.add_run(s["visual_content"])
            if s.get("spoken_content"):
                p = doc.add_paragraph()
                p.add_run("🎤 Lời nói: ").bold = True
                p.add_run(s["spoken_content"])
            if s.get("key_points"):
                p = doc.add_paragraph()
                p.add_run("🔑 Ý chính:").bold = True
                for kp in s["key_points"]:
                    doc.add_paragraph(kp, style="List Bullet")
            if s.get("notes"):
                p = doc.add_paragraph()
                p.add_run("📝 Ghi chú: ").bold = True
                p.add_run(s["notes"])

    exercises = data.get("exercises", [])
    if exercises:
        doc.add_page_break()
        doc.add_heading("Bài tập", level=1)
        for ex in exercises:
            heading_text = ex.get("number", "")
            if ex.get("topic"):
                heading_text += f" — {ex['topic']}"
            doc.add_heading(heading_text, level=2)

            if ex.get("timestamp"):
                p = doc.add_paragraph()
                p.add_run(f"Thời điểm: {ex['timestamp']}").italic = True

            p = doc.add_paragraph()
            p.add_run("Đề bài: ").bold = True
            p.add_run(ex.get("problem_statement", ""))

            if ex.get("given_info"):
                p = doc.add_paragraph()
                p.add_run("Dữ kiện: ").bold = True
                p.add_run(ex["given_info"])
            if ex.get("diagram_description"):
                p = doc.add_paragraph()
                p.add_run("Hình vẽ: ").bold = True
                p.add_run(ex["diagram_description"])
            if ex.get("solution"):
                p = doc.add_paragraph()
                p.add_run("Lời giải: ").bold = True
                p.add_run(ex["solution"])
            if ex.get("answer"):
                p = doc.add_paragraph()
                run = p.add_run("Đáp án: ")
                run.bold = True
                ans = p.add_run(ex["answer"])
                ans.bold = True
                ans.font.color.rgb = RGBColor(0, 120, 0)

    key_terms = data.get("key_terms", [])
    if key_terms:
        doc.add_heading("Thuật ngữ & Định nghĩa", level=1)
        for k in key_terms:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(k.get("term", "")).bold = True
            p.add_run(f": {k.get('definition', '')}")

    if data.get("full_transcript"):
        doc.add_page_break()
        doc.add_heading("Toàn bộ lời nói (Transcript)", level=1)
        for para in data["full_transcript"].split("\n"):
            if para.strip():
                doc.add_paragraph(para)

    doc.save(path)
    return path


# ----------------------- PDF -----------------------
def _escape(text):
    if not text:
        return ""
    return (
        str(text)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace("\n", "<br/>")
    )


def export_pdf(data, output_dir, base_name):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # Register a font that supports Vietnamese diacritics
    font_name = "Helvetica"
    font_bold = "Helvetica-Bold"
    candidate_fonts = []
    if platform.system() == "Windows":
        candidate_fonts = [
            ("VnFont", "C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf"),
            ("VnFont", "C:/Windows/Fonts/tahoma.ttf", "C:/Windows/Fonts/tahomabd.ttf"),
            ("VnFont", "C:/Windows/Fonts/calibri.ttf", "C:/Windows/Fonts/calibrib.ttf"),
        ]
    elif platform.system() == "Darwin":
        candidate_fonts = [
            ("VnFont", "/Library/Fonts/Arial.ttf", "/Library/Fonts/Arial Bold.ttf"),
        ]
    else:
        candidate_fonts = [
            ("VnFont", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
             "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ]
    for name, regular, bold in candidate_fonts:
        try:
            pdfmetrics.registerFont(TTFont(name, regular))
            pdfmetrics.registerFont(TTFont(name + "-Bold", bold))
            font_name = name
            font_bold = name + "-Bold"
            break
        except Exception:
            continue

    path = output_dir / f"{base_name}.pdf"
    pdf = SimpleDocTemplate(
        str(path), pagesize=A4,
        topMargin=2 * cm, bottomMargin=2 * cm,
        leftMargin=2 * cm, rightMargin=2 * cm,
        title=data.get("title", base_name),
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title", parent=styles["Title"],
                                 fontName=font_bold, fontSize=20, alignment=TA_CENTER, spaceAfter=14)
    h1 = ParagraphStyle("H1", parent=styles["Heading1"],
                        fontName=font_bold, fontSize=15, spaceBefore=12, spaceAfter=8, textColor="#1F3864")
    h2 = ParagraphStyle("H2", parent=styles["Heading2"],
                        fontName=font_bold, fontSize=12, spaceBefore=8, spaceAfter=4, textColor="#2E5395")
    body = ParagraphStyle("Body", parent=styles["BodyText"],
                          fontName=font_name, fontSize=10.5, alignment=TA_JUSTIFY, leading=15)
    italic = ParagraphStyle("Italic", parent=body, fontName=font_name, fontSize=9, textColor="#666666")

    story = []
    story.append(Paragraph(_escape(data.get("title", base_name)), title_style))
    story.append(Paragraph(
        f"<i>Loại: {_escape(data.get('content_type', ''))} | "
        f"Ngôn ngữ: {_escape(data.get('language', ''))} | "
        f"Thời lượng: {_escape(data.get('duration_estimate', ''))}</i>",
        italic,
    ))
    story.append(Spacer(1, 0.4 * cm))

    if data.get("summary"):
        story.append(Paragraph("Tóm tắt", h1))
        story.append(Paragraph(_escape(data["summary"]), body))

    sections = data.get("sections", [])
    if sections:
        story.append(Paragraph("Nội dung chi tiết", h1))
        for i, s in enumerate(sections, 1):
            story.append(Paragraph(
                f"{i}. [{_escape(s.get('timestamp_start', ''))} - {_escape(s.get('timestamp_end', ''))}] "
                f"{_escape(s.get('title', ''))}",
                h2,
            ))
            if s.get("visual_content"):
                story.append(Paragraph(f"<b>📺 Hình ảnh:</b> {_escape(s['visual_content'])}", body))
            if s.get("spoken_content"):
                story.append(Paragraph(f"<b>🎤 Lời nói:</b> {_escape(s['spoken_content'])}", body))
            if s.get("notes"):
                story.append(Paragraph(f"<b>📝 Ghi chú:</b> {_escape(s['notes'])}", body))
            story.append(Spacer(1, 0.2 * cm))

    exercises = data.get("exercises", [])
    if exercises:
        story.append(PageBreak())
        story.append(Paragraph("Bài tập", h1))
        for ex in exercises:
            head = _escape(ex.get("number", ""))
            if ex.get("topic"):
                head += f" — {_escape(ex['topic'])}"
            story.append(Paragraph(head, h2))
            if ex.get("timestamp"):
                story.append(Paragraph(f"<i>Thời điểm: {_escape(ex['timestamp'])}</i>", italic))
            story.append(Paragraph(f"<b>Đề bài:</b> {_escape(ex.get('problem_statement', ''))}", body))
            if ex.get("given_info"):
                story.append(Paragraph(f"<b>Dữ kiện:</b> {_escape(ex['given_info'])}", body))
            if ex.get("diagram_description"):
                story.append(Paragraph(f"<b>Hình vẽ:</b> {_escape(ex['diagram_description'])}", body))
            if ex.get("solution"):
                story.append(Paragraph(f"<b>Lời giải:</b> {_escape(ex['solution'])}", body))
            if ex.get("answer"):
                story.append(Paragraph(
                    f"<b>✓ Đáp án:</b> <font color='#008000'>{_escape(ex['answer'])}</font>", body))
            story.append(Spacer(1, 0.3 * cm))

    if data.get("key_terms"):
        story.append(Paragraph("Thuật ngữ & Định nghĩa", h1))
        for k in data["key_terms"]:
            story.append(Paragraph(
                f"• <b>{_escape(k.get('term', ''))}</b>: {_escape(k.get('definition', ''))}", body))

    if data.get("full_transcript"):
        story.append(PageBreak())
        story.append(Paragraph("Toàn bộ lời nói (Transcript)", h1))
        for para in data["full_transcript"].split("\n"):
            if para.strip():
                story.append(Paragraph(_escape(para), body))

    pdf.build(story)
    return path


# ----------------------- XLSX -----------------------
def export_xlsx(data, output_dir, base_name):
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    path = output_dir / f"{base_name}.xlsx"
    wb = Workbook()

    header_fill_blue = PatternFill("solid", fgColor="4472C4")
    header_fill_green = PatternFill("solid", fgColor="70AD47")
    header_fill_orange = PatternFill("solid", fgColor="ED7D31")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    thin = Side(border_style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    wrap_top = Alignment(wrap_text=True, vertical="top")

    # Sheet 1: Overview
    ws = wb.active
    ws.title = "Tổng quan"
    ws["A1"] = "Thuộc tính"
    ws["B1"] = "Giá trị"
    for c in ("A1", "B1"):
        ws[c].font = header_font
        ws[c].fill = header_fill_blue
    rows = [
        ("Tiêu đề", data.get("title", "")),
        ("Loại nội dung", data.get("content_type", "")),
        ("Ngôn ngữ", data.get("language", "")),
        ("Thời lượng ước tính", data.get("duration_estimate", "")),
        ("Tóm tắt", data.get("summary", "")),
        ("Số sections", len(data.get("sections", []))),
        ("Số bài tập", len(data.get("exercises", []))),
        ("Số thuật ngữ", len(data.get("key_terms", []))),
    ]
    for i, (k, v) in enumerate(rows, 2):
        ws.cell(i, 1, k).font = Font(bold=True)
        ws.cell(i, 2, v).alignment = wrap_top
    ws.column_dimensions["A"].width = 24
    ws.column_dimensions["B"].width = 90

    # Sheet 2: Timeline / Sections
    ws2 = wb.create_sheet("Timeline")
    headers = ["STT", "Bắt đầu", "Kết thúc", "Tiêu đề", "Hình ảnh/Slide", "Lời nói", "Ý chính", "Ghi chú"]
    for col, h in enumerate(headers, 1):
        c = ws2.cell(1, col, h)
        c.font = header_font
        c.fill = header_fill_blue
        c.border = border
    for i, s in enumerate(data.get("sections", []), 2):
        ws2.cell(i, 1, i - 1)
        ws2.cell(i, 2, s.get("timestamp_start", ""))
        ws2.cell(i, 3, s.get("timestamp_end", ""))
        ws2.cell(i, 4, s.get("title", ""))
        ws2.cell(i, 5, s.get("visual_content", ""))
        ws2.cell(i, 6, s.get("spoken_content", ""))
        ws2.cell(i, 7, "\n".join(f"• {kp}" for kp in s.get("key_points", []) or []))
        ws2.cell(i, 8, s.get("notes", ""))
        for col in range(1, 9):
            ws2.cell(i, col).alignment = wrap_top
            ws2.cell(i, col).border = border
    widths = [5, 12, 12, 30, 50, 60, 40, 25]
    for i, w in enumerate(widths, 1):
        ws2.column_dimensions[get_column_letter(i)].width = w
    ws2.freeze_panes = "A2"

    # Sheet 3: Exercises
    ws3 = wb.create_sheet("Bài tập")
    headers = ["STT", "Bài số", "Thời điểm", "Chủ đề",
               "Đề bài", "Dữ kiện", "Hình vẽ", "Lời giải", "Đáp án", "Độ khó"]
    for col, h in enumerate(headers, 1):
        c = ws3.cell(1, col, h)
        c.font = header_font
        c.fill = header_fill_green
        c.border = border
    for i, ex in enumerate(data.get("exercises", []), 2):
        ws3.cell(i, 1, i - 1)
        ws3.cell(i, 2, ex.get("number", ""))
        ws3.cell(i, 3, ex.get("timestamp", ""))
        ws3.cell(i, 4, ex.get("topic", ""))
        ws3.cell(i, 5, ex.get("problem_statement", ""))
        ws3.cell(i, 6, ex.get("given_info", ""))
        ws3.cell(i, 7, ex.get("diagram_description", ""))
        ws3.cell(i, 8, ex.get("solution", ""))
        ws3.cell(i, 9, ex.get("answer", ""))
        ws3.cell(i, 10, ex.get("difficulty", ""))
        for col in range(1, 11):
            ws3.cell(i, col).alignment = wrap_top
            ws3.cell(i, col).border = border
    widths = [5, 10, 12, 20, 50, 30, 40, 50, 25, 10]
    for i, w in enumerate(widths, 1):
        ws3.column_dimensions[get_column_letter(i)].width = w
    ws3.freeze_panes = "A2"

    # Sheet 4: Key terms
    ws4 = wb.create_sheet("Thuật ngữ")
    for col, h in enumerate(["Thuật ngữ", "Định nghĩa"], 1):
        c = ws4.cell(1, col, h)
        c.font = header_font
        c.fill = header_fill_orange
    for i, k in enumerate(data.get("key_terms", []), 2):
        ws4.cell(i, 1, k.get("term", "")).font = Font(bold=True)
        ws4.cell(i, 2, k.get("definition", "")).alignment = wrap_top
    ws4.column_dimensions["A"].width = 28
    ws4.column_dimensions["B"].width = 70
    ws4.freeze_panes = "A2"

    # Sheet 5: Full transcript
    if data.get("full_transcript"):
        ws5 = wb.create_sheet("Transcript")
        ws5["A1"] = "Toàn bộ lời nói"
        ws5["A1"].font = header_font
        ws5["A1"].fill = header_fill_blue
        ws5["A2"] = data["full_transcript"]
        ws5["A2"].alignment = wrap_top
        ws5.column_dimensions["A"].width = 120

    wb.save(path)
    return path


# ----------------------- PPTX -----------------------
def export_pptx(data, output_dir, base_name):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = output_dir / f"{base_name}.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return slide

    def add_content_slide(title, body_blocks):
        """body_blocks: list of (label, text, is_bold) tuples"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
        slide.shapes.title.text = title[:120]
        # Add a textbox for content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(12.3)
        height = Inches(5.5)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for label, text, is_bold in body_blocks:
            if not text:
                continue
            text = str(text)
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            if label:
                run = p.add_run()
                run.text = label
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
            run2 = p.add_run()
            run2.text = " " + text[:800] + ("..." if len(text) > 800 else "")
            run2.font.size = Pt(12)
            run2.font.bold = is_bold
        return slide

    # Title slide
    add_title_slide(
        data.get("title", base_name),
        f"{data.get('content_type', '').title()}  |  {data.get('language', 'vi').upper()}  |  {data.get('duration_estimate', '')}",
    )

    # Summary slide
    if data.get("summary"):
        add_content_slide("Tóm tắt", [("", data["summary"], False)])

    # Sections
    for i, s in enumerate(data.get("sections", []), 1):
        blocks = [
            ("⏱", f"{s.get('timestamp_start', '')} - {s.get('timestamp_end', '')}", True),
        ]
        if s.get("visual_content"):
            blocks.append(("📺 Hình ảnh:", s["visual_content"], False))
        if s.get("spoken_content"):
            blocks.append(("🎤 Lời nói:", s["spoken_content"], False))
        if s.get("key_points"):
            blocks.append(("🔑 Ý chính:", "\n• " + "\n• ".join(s["key_points"]), False))
        add_content_slide(f"{i}. {s.get('title', '')}", blocks)

    # Exercises
    for ex in data.get("exercises", []):
        title = ex.get("number", "Bài tập")
        if ex.get("topic"):
            title += f" — {ex['topic']}"
        blocks = [("Đề bài:", ex.get("problem_statement", ""), False)]
        if ex.get("given_info"):
            blocks.append(("Dữ kiện:", ex["given_info"], False))
        if ex.get("diagram_description"):
            blocks.append(("Hình vẽ:", ex["diagram_description"], False))
        if ex.get("solution"):
            blocks.append(("Lời giải:", ex["solution"], False))
        if ex.get("answer"):
            blocks.append(("✓ Đáp án:", ex["answer"], True))
        add_content_slide(title, blocks)

    # Key terms (one slide)
    if data.get("key_terms"):
        body = "\n".join(f"• {k.get('term', '')}: {k.get('definition', '')}"
                         for k in data["key_terms"])
        add_content_slide("Thuật ngữ & Định nghĩa", [("", body, False)])

    prs.save(path)
    return path
